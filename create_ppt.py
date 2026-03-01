
# I'll create a Python script that generates your PPT with the theme
# This code you can run locally to create the .pptx file

ppt_code = '''
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color definitions (Neon Traffic Tech Theme)
DARK_BG = RGBColor(15, 23, 42)      # #0f172a
CARD_BG = RGBColor(30, 41, 59)      # #1e293b
AMBER = RGBColor(245, 158, 11)      # #f59e0b
GREEN = RGBColor(16, 185, 129)      # #10b981
RED = RGBColor(239, 68, 68)         # #ef4444
CYAN = RGBColor(6, 182, 212)        # #06b6d4
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(148, 163, 184)      # #94a3b8

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = AMBER
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_lines, accent_color=AMBER):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Header bar
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = CARD_BG
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = accent_color
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = line
        p.font.size = Pt(20) if line.startswith('•') or line.startswith('→') else Pt(22)
        p.font.color.rgb = WHITE if line.startswith('•') or line.startswith('→') else GRAY
        p.space_before = Pt(12)
        p.level = 1 if line.startswith('  ') else 0
    
    return slide

# Slide 1: Title
add_title_slide(prs, 
    "🚦 REAL-TIME TRAFFIC CONGESTION PREDICTION", 
    "& RISK INTELLIGENCE SYSTEM\\n\\nPython • Machine Learning • Real-Time APIs")

# Slide 2: Problem Statement
add_content_slide(prs, "THE URBAN TRAFFIC CRISIS", [
    "❌ EXISTING SYSTEMS:",
    "   • Only display current traffic speed",
    "   • No intelligent prediction capability", 
    "   • No risk evaluation or advisory",
    "   • Reactive rather than proactive",
    "",
    "📍 INDIAN CONTEXT:",
    "   • Major cities face 40-50% congestion daily",
    "   • Loss of ₹1.5 lakh crore annually due to traffic",
    "   • No smart decision support for commuters",
    "",
    "💡 OPPORTUNITY:",
    "   Convert raw traffic data into INTELLIGENT ROUTE DECISIONS"
], RED)

# Slide 3: Objectives
add_content_slide(prs, "🎯 PROJECT OBJECTIVES", [
    "🔄 LIVE DATA FETCHING",
    "   • Real-time traffic via TomTom API",
    "",
    "🤖 ML PREDICTION & CLASSIFICATION", 
    "   • Random Forest Classifier (High/Medium/Low)",
    "",
    "⚠️ RISK SCORING ENGINE",
    "   • Calculate Route Risk Score (0-100)",
    "",
    "🧠 SMART TRAVEL ADVISORY",
    "   • Generate real-time travel recommendations",
    "",
    "🏙️ MULTI-CITY COMPARISON",
    "   • Compare multiple cities simultaneously"
], AMBER)

# Slide 4: Architecture
add_content_slide(prs, "🏗️ SYSTEM ARCHITECTURE", [
    "┌─────────────────┐",
    "│   USER INPUT    │  ← City Name (Streamlit)",
    "└────────┬────────┘",
    "         ↓",
    "┌─────────────────┐",
    "│  GEOCODING API  │  ← Nominatim API (Lat/Long)",
    "└────────┬────────┘",
    "         ↓",
    "┌─────────────────┐", 
    "│   TOMTOM API    │  ← Live Speed Data",
    "└────────┬────────┘",
    "         ↓",
    "┌─────────────────┐",
    "│  ML PREDICTION  │  ← Random Forest Classifier",
    "│   RISK ENGINE   │",
    "└────────┬────────┘",
    "         ↓",
    "┌─────────────────┐",
    "│   DASHBOARD     │  ← Streamlit + Plotly + Folium",
    "│ VISUALIZATION   │",
    "└─────────────────┘"
], CYAN)

# Slide 5: Technology Stack
add_content_slide(prs, "🛠️ TECHNOLOGY STACK", [
    "PROGRAMMING: 🐍 Python (Core Logic)",
    "",
    "DASHBOARD: 🌐 Streamlit (Interactive Web UI)",
    "",
    "MAPS: 🗺️ Folium (Interactive Maps)",
    "",
    "VISUALIZATION: 📊 Plotly (Charts & Gauge Meter)",
    "",
    "MACHINE LEARNING: 🤖 Scikit-learn (Random Forest)",
    "",
    "DATA HANDLING: 🐼 Pandas",
    "",
    "API CALLS: 🌐 Requests",
    "",
    "EXTERNAL APIs:",
    "   • TomTom Traffic API (Real-time data)",
    "   • Nominatim Geocoding API (Location)"
], GREEN)

# Slide 6: Key Features 1
add_content_slide(prs, "⚡ KEY FEATURES - I", [
    "🚦 LIVE TRAFFIC DATA FETCHING",
    "   • Real-time speed & free flow speed monitoring",
    "   • 30-second auto-refresh capability",
    "",
    "🤖 ML-BASED CONGESTION PREDICTION",
    "   • Algorithm: Random Forest Classifier",
    "   • Input: Speed ratio + Time features",
    "   • Output: High / Medium / Low classification",
    "",
    "📊 CONGESTION GAUGE METER",
    "   • Visual 0-100% congestion display",
    "   • 🟢 0-40% (Low) | 🟡 40-70% (Medium) | 🔴 70-100% (High)",
    "   • Real-time animated needle"
], AMBER)

# Slide 7: Key Features 2
add_content_slide(prs, "⚡ KEY FEATURES - II", [
    "⏰ PEAK HOUR DETECTION",
    "   • Morning Rush: 8:00 AM - 11:00 AM",
    "   • Evening Rush: 5:00 PM - 9:00 PM", 
    "   • Auto-adjusts risk weightage",
    "",
    "🚨 ROUTE RISK SCORE ENGINE (0-100)",
    "   Formula: Congestion% + Peak Penalty + Speed Penalty",
    "   • >70: AVOID | 40-70: CAUTION | <40: SAFE",
    "",
    "🏙️ MULTI-CITY COMPARISON",
    "   • Compare multiple cities simultaneously",
    "   • Risk ranking with medals 🥇🥈🥉",
    "   • Smart recommendation engine"
], RED)

# Slide 8: ML Model
add_content_slide(prs, "🤖 MACHINE LEARNING MODEL", [
    "ALGORITHM: Random Forest Classifier",
    "",
    "WHY RANDOM FOREST?",
    "   ✅ Handles non-linear traffic patterns",
    "   ✅ Reduces overfitting (ensemble learning)",
    "   ✅ Excellent for classification problems",
    "   ✅ Robust with small datasets",
    "",
    "TRAINING FEATURES:",
    "   • Current Speed (km/h) | Free Flow Speed (km/h)",
    "   • Hour of Day (0-23) | Speed Ratio (Current/Free)",
    "",
    "OUTPUT CLASSES:",
    "   🟢 Low Congestion | 🟡 Medium Congestion | 🔴 High Congestion",
    "",
    "PERFORMANCE: ~85-90% Accuracy | <100ms Prediction Time"
], CYAN)

# Slide 9: Risk Score Logic
add_content_slide(prs, "🧮 RISK SCORE CALCULATION", [
    "FORMULA:",
    "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
    "RISK SCORE = Base Congestion %",
    "           + Peak Hour Penalty (+10)",
    "           + Speed Ratio Penalty (+10)",
    "",
    "FINAL RANGE: 0 (Safe) — 100 (Avoid)",
    "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
    "",
    "EXAMPLE:",
    "   Current Speed: 25 km/h | Free Flow: 50 km/h",
    "   Congestion: 50% | Peak Hour: YES | Speed Ratio: 0.5",
    "",
    "   Calculation: 50% + 10 + 10 = 70 RISK SCORE",
    "",
    "   RESULT: 🟡 CAUTION - Heavy Traffic Expected"
], AMBER)

# Slide 10: Smart Advisory
add_content_slide(prs, "🧠 SMART TRAVEL ADVISORY", [
    "DECISION TABLE:",
    "",
    "┌──────────────┬─────────────┬─────────────────────┐",
    "│ ML PREDICTION│ PEAK HOUR   │ ADVISORY            │",
    "├──────────────┼─────────────┼─────────────────────┤",
    "│ 🔴 HIGH      │ YES         │ ❌ AVOID THIS ROUTE │",
    "│ 🔴 HIGH      │ NO          │ ⚠️ HEAVY TRAFFIC    │",
    "│ 🟡 MEDIUM    │ YES         │ 🚗 TRAVEL CAREFULLY │",
    "│ 🟢 LOW       │ ANY         │ ✅ SAFE TO TRAVEL   │",
    "└──────────────┴─────────────┴─────────────────────┘",
    "",
    "💡 DYNAMIC: Advisory refreshes every 30 seconds",
    "   based on live API data and ML predictions"
], GREEN)

# Slide 11: Multi-City
add_content_slide(prs, "🏙️ MULTI-CITY COMPARISON", [
    "LIVE COMPARISON DASHBOARD:",
    "",
    "🥇 Mumbai    Speed: 22 km/h  Risk: 85  [🔴 HIGH]",
    "🥈 Delhi     Speed: 35 km/h  Risk: 60  [🟡 MED]",
    "🥉 Bangalore Speed: 42 km/h  Risk: 45  [🟡 MED]",
    "4️⃣ Pune      Speed: 50 km/h  Risk: 30  [🟢 LOW]",
    "",
    "RISK SUMMARY: High: 1 | Medium: 2 | Low: 1",
    "",
    "SMART RECOMMENDATION:",
    "   ❌ \\"Avoid Mumbai - Highest congestion detected\\"",
    "   ✅ \\"Safest choice: Pune - Low traffic risk\\"",
    "",
    "📊 Visual: Risk comparison bar chart (Plotly)"
], CYAN)

# Slide 12: Innovation
add_content_slide(prs, "💡 PROJECT INNOVATION", [
    "WHAT MAKES IT UNIQUE?",
    "",
    "🔗 INTEGRATION:",
    "   • Real-time APIs + Machine Learning + Risk Analytics",
    "",
    "🧮 CUSTOM RISK ENGINE:",
    "   • Composite scoring (0-100) with multiple factors",
    "",
    "🏆 MULTI-CITY RANKING:",
    "   • Comparative analytics with medal system",
    "",
    "🧠 INTELLIGENT ADVISORY:",
    "   • Context-aware recommendations (not just data display)",
    "",
    "IMPACT: Transforms raw data into actionable decisions",
    "        for smarter urban mobility"
], AMBER)

# Slide 13: Future Scope
add_content_slide(prs, "🚀 FUTURE SCOPE", [
    "ENHANCEMENTS ROADMAP:",
    "",
    "🗺️ INDIA TRAFFIC HEATMAP",
    "   • Visual heatmap of congestion across cities",
    "",
    "🤖 DEEP LEARNING (LSTM)",
    "   • Time-series prediction for traffic forecasting",
    "",
    "☁️ CLOUD DEPLOYMENT",
    "   • AWS/GCP hosting for 24/7 availability",
    "",
    "🗺️ ROUTE OPTIMIZATION",
    "   • AI-powered alternative route suggestions",
    "",
    "📱 MOBILE APPLICATION",
    "   • iOS/Android app for on-the-go access",
    "",
    "🔔 PUSH NOTIFICATIONS",
    "   • Alerts for high-risk routes before travel"
], GREEN)

# Slide 14: Conclusion
add_content_slide(prs, "✅ CONCLUSION", [
    "🎯 TRANSFORMED raw traffic data into",
    "   INTELLIGENT ROUTE DECISIONS",
    "",
    "🔧 INTEGRATED multiple technologies:",
    "   • Real-time APIs • Machine Learning",
    "   • Risk Analytics • Interactive Dashboard",
    "",
    "💡 DELIVERED actionable insights:",
    "   • Risk scores (0-100) • Smart advisories",
    "   • Multi-city comparison • Peak hour detection",
    "",
    "🚀 IMPACT: Helps commuters make informed decisions,",
    "    saving time & reducing urban congestion stress",
    "",
    "   \\"From Data to Decisions\\" 🚦➡️🧠"
], CYAN)

# Slide 15: Thank You
add_title_slide(prs, "🙏 THANK YOU", "QUESTIONS & ANSWERS\\n\\n📧 [your.email@example.com] | 🔗 [GitHub Link]")

# Save
prs.save('Traffic_Congestion_Prediction_PPT.pptx')
print("✅ PPT Created Successfully!")
print("📁 File: Traffic_Congestion_Prediction_PPT.pptx")
print("🎨 Theme: Neon Traffic Tech (Dark Blue + Amber)")
'''

print("="*60)
print("🐍 PYTHON CODE TO GENERATE YOUR PPT")
print("="*60)
print(ppt_code)
print("="*60)
print("\n✅ INSTRUCTIONS:")
print("1. Install library: pip install python-pptx")
print("2. Save above code as 'create_ppt.py'")
print("3. Run: python create_ppt.py")
print("4. Your PPT will be created with the Neon Traffic Tech theme!")

# ================================
# EXECUTE THE CODE TO CREATE PPT
# ================================

exec(ppt_code)